"""
Document preview converter — extracts text content for non-PDF files
"""
import os
import json


def convert_to_preview(file_path, file_type):
    """Convert document to previewable format"""

    file_type = (file_type or '').lower()

    try:
        # Plain text files
        if file_type in ('txt', 'md', 'log', 'csv', 'json', 'xml', 'html', 'htm', 'rtf'):
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return {
                'type': 'text',
                'content': content[:50000],  # limit to 50k chars
                'format': file_type,
            }

        # DOCX
        if file_type == 'docx':
            try:
                import mammoth
                with open(file_path, 'rb') as f:
                    result = mammoth.convert_to_html(f)
                return {
                    'type': 'html',
                    'content': result.value,
                    'format': 'docx',
                }
            except ImportError:
                return {
                    'type': 'text',
                    'content': '[Install mammoth package: pip install mammoth]',
                    'format': 'docx',
                }

        # PPTX
        if file_type == 'pptx':
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                slides = []
                for i, slide in enumerate(prs.slides):
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            texts.append(shape.text_frame.text)
                    slides.append({
                        'number': i + 1,
                        'content': '\n'.join(texts),
                    })
                return {
                    'type': 'slides',
                    'slides': slides,
                    'total_slides': len(slides),
                    'format': 'pptx',
                }
            except ImportError:
                return {
                    'type': 'text',
                    'content': '[Install python-pptx package: pip install python-pptx]',
                    'format': 'pptx',
                }

        # PDF — handled by frontend directly
        if file_type == 'pdf':
            return {
                'type': 'pdf',
                'content': 'Use /api/documents/{id}/file endpoint',
                'format': 'pdf',
            }

        return {
            'type': 'unsupported',
            'content': f'Preview not available for .{file_type} files',
            'format': file_type,
        }

    except Exception as e:
        return {
            'type': 'error',
            'content': f'Error generating preview: {str(e)}',
            'format': file_type,
        }